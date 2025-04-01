import os
import streamlit as st
from langchain.chains import RetrievalQA
from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadError
import pdfplumber
from PIL import Image
import docx
from pptx import Presentation
import easyocr
import fitz 
from langchain.callbacks.base import BaseCallbackHandler
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.chroma import Chroma
from langchain_community.vectorstores import Neo4jVector
from streamlit.logger import get_logger
from langchain_community.document_loaders.pdf import PyPDFLoader
from langchain_community.document_loaders import CSVLoader
import pandas as pd
from chains import load_embedding_model, load_llm
from utils.file_readers import read_pdf, read_csv, read_excel, read_msg
from dotenv import load_dotenv
from neo4j import GraphDatabase
import io
import extract_msg
import shutil
import zipfile
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
from langchain_openai import OpenAIEmbeddings
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.embeddings import BedrockEmbeddings
from langchain_community.embeddings.sentence_transformer import SentenceTransformerEmbeddings

from langchain_openai import ChatOpenAI
from langchain_community.chat_models import ChatOllama
from langchain_community.chat_models import BedrockChat
from langchain_google_genai import GoogleGenerativeAIEmbeddings

load_dotenv(".env")

url = os.getenv("NEO4J_URI")
username = os.getenv("NEO4J_USERNAME")
password = os.getenv("NEO4J_PASSWORD")
ollama_base_url = os.getenv("OLLAMA_BASE_URL")
embedding_model_name = os.getenv("EMBEDDING_MODEL")
llm_name = os.getenv("LLM")

class BaseLogger:
    def __init__(self) -> None:
        self.info = print


def extract_title_and_question(input_string):
    lines = input_string.strip().split("\n")

    title = ""
    question = ""
    is_question = False  # flag to know if we are inside a "Question" block

    for line in lines:
        if line.startswith("Title:"):
            title = line.split("Title: ", 1)[1].strip()
        elif line.startswith("Question:"):
            question = line.split("Question: ", 1)[1].strip()
            is_question = (
                True  # set the flag to True once we encounter a "Question:" line
            )
        elif is_question:
            # if the line does not start with "Question:" but we are inside a "Question" block,
            # then it is a continuation of the question
            question += "\n" + line.strip()

    return title, question

def load_embedding_model(embedding_model_name: str, logger=BaseLogger(), config={}):
    if embedding_model_name == "ollama":
        embeddings = OllamaEmbeddings(
            base_url=config["ollama_base_url"], model="llama3-70b-instruct-v1:0"
        )
        dimension = 4096
        logger.info("Embedding: Using Ollama")
    elif embedding_model_name == "openai":
        embeddings = OpenAIEmbeddings()
        dimension = 1536
        logger.info("Embedding: Using OpenAI")
    elif embedding_model_name == "aws":
        embeddings = BedrockEmbeddings()
        dimension = 1536
        logger.info("Embedding: Using AWS")
    elif embedding_model_name == "google-genai-embedding-001":        
        embeddings = GoogleGenerativeAIEmbeddings(
            model="models/embedding-001"
        )
        dimension = 768
        logger.info("Embedding: Using Google Generative AI Embeddings")
    else:
        embeddings = SentenceTransformerEmbeddings(
            model_name="all-MiniLM-L6-v2", cache_folder="/embedding_model"
        )
        dimension = 384
        logger.info("Embedding: Using SentenceTransformer")
    return embeddings, dimension

def load_llm(llm_name: str, logger=BaseLogger(), config={}):
    if llm_name == "gpt-4":
        logger.info("LLM: Using GPT-4")
        return ChatOpenAI(temperature=0, model_name="gpt-4", streaming=True)
    elif llm_name == "gpt-3.5":
        logger.info("LLM: Using GPT-3.5")
        return ChatOpenAI(temperature=0, model_name="gpt-3.5-turbo", streaming=True)
    elif llm_name == "claudev2":
        logger.info("LLM: ClaudeV2")
        return BedrockChat(
            model_id="anthropic.claude-3-5-sonnet-20240620-v1:0",
            model_kwargs={"temperature": 0.0},
            streaming=True
        )
    elif len(llm_name):
        logger.info(f"LLM: Using Ollama: {llm_name}")
        return ChatOllama(
            temperature=0,
            base_url=config["ollama_base_url"],
            model=llm_name,
            streaming=True,
            # seed=2,
            top_k=10,  # A higher value (100) will give more diverse answers, while a lower value (10) will be more conservative.
            top_p=0.3,  # Higher value (0.95) will lead to more diverse text, while a lower value (0.5) will generate more focused text.
            num_ctx=3072  # Sets the size of the context window used to generate the next token.
        )
    logger.info("LLM: Using GPT-3.5")
    return ChatOpenAI(temperature=0, model_name="gpt-3.5-turbo", streaming=True)

os.environ["NEO4J_URL"] = url

logger = get_logger(__name__)

embeddings, dimension = load_embedding_model(
    embedding_model_name, config={"ollama_base_url": ollama_base_url}, logger=logger
)


class DocumentLoader: # For .txt files
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        doc = docx.Document(self.file_path)
        line_number = 0
        page_number = 1  # Page number starts from 1

        for paragraph in doc.paragraphs:

            yield Document(
                page_content=paragraph.text,
                metadata={"line_number": line_number, "page_number": page_number, "source": self.file_path},
            )
            line_number += 1


class PptxLoader: # For .ppt or .pptx files
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        prs = Presentation(self.file_path)
        slide_number = 0

        for slide in prs.slides:
            slide_number += 1
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    yield Document(
                        page_content=shape.text,
                        metadata={"slide_number": slide_number, "source": self.file_path},
                    )


def create_vectorstore(index_name, chunks):
    driver = GraphDatabase.driver(url, auth=(username, password))
    with driver.session() as session:
        # Create index if it doesn't exist
        session.run(f"CREATE INDEX {index_name} IF NOT EXISTS FOR (n:Chunk) ON (n.embedding)")
        

    vectorstore = Chroma.from_texts(texts=chunks,
                                    embedding=embeddings)
    
    # Neo4jVector.from_texts(
    #     chunks,
    #     url=url,
    #     username=username,
    #     password=password,
    #     embedding=embeddings,
    #     index_name=index_name,
    #     node_label="Chunk",
    #     pre_delete_collection=False,
    # )
    
    return vectorstore

class StreamHandler(BaseCallbackHandler):
    def __init__(self, container, initial_text=""):
        self.container = container
        self.text = initial_text

    def on_llm_new_token(self, token: str, **kwargs) -> None:
        self.text += token
        self.container.markdown(self.text)

llm = load_llm(llm_name, logger=logger, config={"ollama_base_url": ollama_base_url})

# Global variable to hold the last response
r1 = ""

# Function to parse the last_response into a structured DataFrame
def parse_response(response):
    # Assuming the response is in markdown table format, split by newline
    lines = response.split('\n')

    # Find the header line and the start of the data
    header_line = None
    data_start = None
    for i, line in enumerate(lines):
        if '---' in line:
            header_line = lines[i - 1]  # The line before the separator line is the header
            data_start = i + 1  # Data starts after the separator line
            break

    if header_line is None or data_start is None:
        return pd.DataFrame()  # Return an empty DataFrame if headers or data start are not found
    
    # Extract headers and data
    headers = [header.strip() for header in header_line.split('|')]
    data = []
    for line in lines[data_start:]:
        columns = line.split('|')
        if len(columns) == len(headers):
            data.append({headers[i]: columns[i].strip() for i in range(len(headers))})
        else:
            # Handle rows with missing or extra columns
            row_data = {headers[i]: (columns[i].strip() if i < len(columns) else '') for i in range(len(headers))}
            data.append(row_data)

    # Creating DataFrame from the list of dictionaries
    df = pd.DataFrame(data)
    return df

def check_embedding_similarity(embeddings, threshold=0.95):
    """
    Checks for near-duplicate embeddings based on cosine similarity.
    Parameters:
        - embeddings: List of embeddings to compare.
        - threshold: Similarity threshold above which items are considered duplicates.
    """
    # Compute pairwise cosine similarity matrix
    similarity_matrix = cosine_similarity(embeddings)

    # Check each pair of embeddings for high similarity
    n = len(embeddings)
    duplicate_found = False
    for i in range(n):
        for j in range(i + 1, n):
            if similarity_matrix[i, j] > threshold:
                st.write(f"High similarity ({similarity_matrix[i, j]:.2f}) between embeddings {i + 1} and {j + 1}")
                duplicate_found = True
    if not duplicate_found:
        st.write("No near-duplicate embeddings found.")

def create_or_update_vectorstore(index_name, new_chunks):
    # Check if the vectorstore already exists in session state
    if index_name in st.session_state:
        existing_chunks = st.session_state[index_name]
        # Remove duplicates before updating the vector store
        unique_chunks = list(set(existing_chunks + new_chunks))
        st.session_state[index_name].extend(unique_chunks)
    else:
        # If it doesn't exist, create a new vectorstore
        unique_chunks = list(set(new_chunks))
        st.session_state[index_name] = unique_chunks

    # Generate embeddings for each chunk and display them
    #chunk_embeddings = []
    #for i, chunk in enumerate(st.session_state[index_name]):
        #embedding = embeddings.embed_documents([chunk])[0]
        #chunk_embeddings.append(embedding)
        #st.write(f"Embedding for Chunk {i + 1}:", embedding)  # Display the embedding

    # Check for near-duplicate embeddings
    #check_embedding_similarity(chunk_embeddings)

    # Create or update the Chroma vectorstore with the updated chunks
    vectorstore = Chroma.from_texts(texts=st.session_state[index_name], embedding=embeddings)
    #st.write(f"Vectorstore created with {len(st.session_state[index_name])} chunks")
    return vectorstore


def main():

    # Initialize session state variables if they do not exist
    if 'messages' not in st.session_state:
        st.session_state.messages = []
        
    if 'pdf_chunks' not in st.session_state:
        st.session_state.pdf_chunks = []

    if 'csv_chunks' not in st.session_state:
        st.session_state.csv_chunks = []

    if 'excel_chunks' not in st.session_state:
        st.session_state.excel_chunks = []

    if 'vectorstore' not in st.session_state:
        st.session_state.vectorstore = None
       
    global r1
    
    st.title("📄 Chat with Your File")

    if 'messages' not in st.session_state:
        st.session_state.messages = []

    # Sidebar for chat history
    st.sidebar.title("Chat History")
    for i, msg in enumerate(st.session_state.messages):
        st.sidebar.write(f"{i + 1}. {msg['sender']}: {msg['text']}")

    # Main content area
    st.header("Upload and Chat with your files")

    uploaded_files = st.file_uploader("Upload your files1", type=["pdf", "csv", "xlsx"], accept_multiple_files=True)
    
    if uploaded_files is not None:
        
        for uploaded_file in uploaded_files:
            
            if uploaded_file.name.endswith(".pdf"):
                
                # Read the PDF and get both text and image data
                pdf_content = read_pdf(uploaded_file)

                text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50)
                chunks = text_splitter.split_text(text=pdf_content)
                
                # Accumulate chunks for PDFs in session state
                st.session_state.pdf_chunks.extend(chunks)
                
            elif uploaded_file.name.endswith(".csv"):
                df = read_csv(uploaded_file)
                chunks = df.apply(lambda row: row.dropna().to_string(index=False), axis=1).tolist()
                #st.write("CSV chunks:", chunks)  # Check the chunks generated
                #st.write("Number of CSV chunks:", len(chunks))  # Check the number of chunks

                # Accumulate chunks for CSVs in session state
                st.session_state.csv_chunks.extend(chunks)

            elif uploaded_file.name.endswith(".xlsx"):
                df = read_excel(uploaded_file)
                chunks = df.apply(lambda row: row.to_string(index=False), axis=1).tolist()

                # Accumulate chunks for Excel files in session state
                st.session_state.excel_chunks.extend(chunks)

            # After processing all files, update vector stores once per file type
            if st.session_state.pdf_chunks:
                st.session_state.vectorstore = create_or_update_vectorstore("pdf_bot", st.session_state.pdf_chunks)
        
            if st.session_state.csv_chunks:
                st.session_state.vectorstore = create_or_update_vectorstore("csv_bot", st.session_state.csv_chunks)

            if st.session_state.excel_chunks:
                st.session_state.vectorstore = create_or_update_vectorstore("excel_bot", st.session_state.excel_chunks)

            # Set up the QA system with the combined vector store
            qa = RetrievalQA.from_chain_type(
            llm=llm, chain_type="stuff", retriever=st.session_state.vectorstore.as_retriever()
            )

        query = st.text_input("Ask questions about your file")
        

        if st.button("Send"):
            if query:
                st.session_state.messages.append({'sender': 'user', 'text': query})
                
                # Construct the conversation history to send to the LLM
                #conversation = "\n".join([f"{msg['sender']}: {msg['text']}" for msg in st.session_state.messages])
                conversation = "".join([f"{msg['sender']}: {msg['text']}" for msg in st.session_state.messages])
                prompt_context = f"You are an expert assistant extracting Pricing from context provided. Provide answer in user defined table format. Context: {conversation} Question: {query} Answer:"

                ###################################################################################################################################

                # Send the entire conversation history to the LLM
                stream_handler = StreamHandler(st.empty())
                response = qa.run(prompt_context, callbacks=[stream_handler])

                st.session_state.messages.append({'sender': 'bot', 'text': response})
                r1 = response  # Save the last response

        if st.button("Test Retrieval"):
            # Ensure the vector store has been created and is stored in session state
            if 'vectorstore' in st.session_state and st.session_state.vectorstore is not None:
                # Get the retriever from the vector store
                retriever = st.session_state.vectorstore.as_retriever(search_type="similarity", k=10)
                
                # Define the test query
                test_query = "What is the price for Eastman DIBK bulk shipments?"
                
                # Retrieve relevant chunks/documents
                relevant_chunks = retriever.get_relevant_documents(test_query)
            
                # Display the retrieved relevant chunks
                if relevant_chunks:
                    for i, chunk in enumerate(relevant_chunks):
                        st.write(f"Relevant Chunk {i + 1}:", chunk.page_content)

                else:
                    st.write("No relevant chunks found.")
            else:
                st.write("Vectorstore is not initialized.")


        if r1:
            response_df = parse_response(r1)

            # Save to CSV (without header)
            csv_buffer = io.StringIO()
            response_df.to_csv(csv_buffer, index=False)
            csv = csv_buffer.getvalue()

            # Show download buttons for CSV and Excel formats

            st.download_button(
                label="Download data as CSV",
                data=csv,
                file_name='selected_df.csv',
                mime='text/csv',
            )
            
            # Provide option to download as Excel
            excel_buffer = io.BytesIO()
            with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
                response_df.to_excel(writer, index=False, sheet_name='Sheet1')
            excel_buffer.seek(0)
            excel_data = excel_buffer.getvalue()

            st.download_button(
                label="Download data as Excel",
                data=excel_data,
                file_name='last_response.xlsx',
                mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            )

    # Separate file uploader for .msg files
    #uploaded_msg_files = []
    extracted_files = []
    

    uploaded_msg_files = st.file_uploader("Upload your Email files", type=["msg"], accept_multiple_files=True)
    
    if uploaded_msg_files is not None:
        for uploaded_msg_file in uploaded_msg_files:
            msg_text, attachments = read_msg(uploaded_msg_file)
            
            # Add the message text and attachments to the list of files to be zipped
            extracted_files.extend([msg_text] + attachments)

    if extracted_files:
        zip_path = 'extracted_files.zip'
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for file in extracted_files:
                with open(file, 'rb') as f:
                    #file_data = f.read()
                    #zipf.writestr(os.path.basename(file), file_data)
                    zipf.write(file, arcname=os.path.basename(file))
                    

        with open(zip_path, 'rb') as zipf:
            st.download_button(
                label="Download all extracted files",
                data=zipf,
                file_name='extracted_files.zip',
                mime='application/zip'
            )

    # Display the conversation
    for msg in st.session_state.messages:
        st.write(f"**{msg['sender'].capitalize()}:** {msg['text']}")

if __name__ == "__main__":
    main()